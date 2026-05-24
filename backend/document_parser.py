import io
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from charset_normalizer import from_bytes


CODE_EXTENSIONS = {
    ".py",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".html",
    ".css",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".sql",
    ".sh",
    ".bash",
    ".go",
    ".rs",
    ".php",
    ".rb",
}

TEXT_EXTENSIONS = {".txt", ".md", ".markdown"}

CODE_LANGUAGE_MAP = {
    ".py": "Python",
    ".java": "Java",
    ".c": "C",
    ".cpp": "C++",
    ".h": "C",
    ".hpp": "C++",
    ".js": "JavaScript",
    ".jsx": "JavaScript",
    ".ts": "TypeScript",
    ".tsx": "TypeScript",
    ".html": "HTML",
    ".css": "CSS",
    ".json": "JSON",
    ".xml": "XML",
    ".yaml": "YAML",
    ".yml": "YAML",
    ".sql": "SQL",
    ".sh": "Shell",
    ".bash": "Bash",
    ".go": "Go",
    ".rs": "Rust",
    ".php": "PHP",
    ".rb": "Ruby",
}

LEGACY_EXTENSIONS = {
    ".doc": "暂不支持旧版 .doc，请先另存为 .docx 后上传。",
    ".ppt": "暂不支持旧版 .ppt，请先另存为 .pptx 后上传。",
}

MAX_NEW_TYPE_SIZE = 20 * 1024 * 1024


def get_file_extension(filename: str) -> str:
    return Path(filename or "").suffix.lower()


def detect_material_type(filename: str, content_type: str | None = None) -> str:
    ext = get_file_extension(filename)
    ctype = (content_type or "").lower()

    if ext == ".pdf" or "pdf" in ctype:
        return "PDF"

    if ext in {".png", ".jpg", ".jpeg", ".webp"} or any(
        tag in ctype for tag in ("image/png", "image/jpeg", "image/webp")
    ):
        return "IMAGE"

    if ext == ".docx":
        return "DOCX"

    if ext == ".pptx":
        return "PPTX"

    if ext in LEGACY_EXTENSIONS:
        return "LEGACY"

    if ext in TEXT_EXTENSIONS:
        return "TEXT"

    if ext in CODE_EXTENSIONS:
        return "CODE"

    return "UNSUPPORTED"


def _read_docx_paragraphs(doc: DocxDocument) -> list[str]:
    lines: list[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style_name = (para.style.name if para.style and para.style.name else "").lower()
        if any(tag in style_name for tag in ("heading", "title", "heading1", "heading2", "heading3")):
            lines.append(f"## {text}")
        else:
            lines.append(text)
    return lines


def _read_docx_tables(doc: DocxDocument) -> list[str]:
    lines: list[str] = []
    for table in doc.tables:
        lines.append("")
        for row_index, row in enumerate(table.rows):
            cells = [(cell.text or "").strip() for cell in row.cells]
            lines.append(" | ".join(cells))
            if row_index == 0:
                lines.append(" | ".join("---" for _ in cells))
        lines.append("")
    return lines


def extract_docx_text(file_bytes: bytes) -> str:
    stream = io.BytesIO(file_bytes)
    try:
        doc = DocxDocument(stream)
    except Exception as e:
        raise ValueError(f"无法解析 Word 文档：{e}")

    parts: list[str] = []
    parts.extend(_read_docx_paragraphs(doc))
    parts.extend(_read_docx_tables(doc))

    result = "\n\n".join(parts).strip()
    if not result:
        result = "（此 Word 文档未提取到可读文本）"
    return result


def extract_pptx_text(file_bytes: bytes) -> str:
    stream = io.BytesIO(file_bytes)
    try:
        prs = Presentation(stream)
    except Exception as e:
        raise ValueError(f"无法解析 PPT 文稿：{e}")

    slides_output: list[str] = []
    total_text = ""

    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"【第 {idx} 页】"]

        title_texts: list[str] = []
        body_texts: list[str] = []
        table_texts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame:
                frame_text = shape.text_frame.text.strip()
                if not frame_text:
                    continue
                if shape.is_placeholder and shape.placeholder_format.type == 1:
                    title_texts.append(frame_text)
                else:
                    body_texts.append(frame_text)

            if shape.has_table and shape.table:
                rows = []
                for row in shape.table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    table_texts.append("\n".join(rows))

        if title_texts:
            slide_parts.append(f"标题：{'；'.join(title_texts)}")
        if body_texts:
            slide_parts.append(f"内容：{'；'.join(body_texts)}")
        if table_texts:
            slide_parts.append(f"表格：{' | '.join(table_texts)}")

        if slide.has_notes_slide and slide.notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"备注：{notes_text}")

        slide_text = "\n".join(slide_parts)
        total_text += slide_text
        slides_output.append(slide_text)

    result = "\n\n".join(slides_output).strip()
    if not result or len(total_text) < 10:
        result = "（此 PPT 未提取到可读文本，可能为纯图片扫描件）"
    return result


def extract_text_file(file_bytes: bytes, filename: str) -> str:
    try:
        detected = from_bytes(file_bytes).best()
        if detected:
            return str(detected)
    except Exception:
        pass

    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue

    return file_bytes.decode("utf-8", errors="replace")


def extract_code_file(file_bytes: bytes, filename: str) -> str:
    ext = get_file_extension(filename)
    language = CODE_LANGUAGE_MAP.get(ext, "Code")
    raw_text = extract_text_file(file_bytes, filename)

    max_len = 500_000
    if len(raw_text) > max_len:
        raw_text = raw_text[:max_len] + "\n\n...（文件过长，已截断）"

    header = f"文件名：{Path(filename).name}\n语言：{language}\n以下是该代码文件的源码：\n"
    return header + raw_text


def extract_supported_file_text(
    file_bytes: bytes, filename: str, content_type: str | None = None
) -> dict:
    material_type = detect_material_type(filename, content_type)
    ext = get_file_extension(filename)

    if material_type == "LEGACY":
        raise ValueError(LEGACY_EXTENSIONS.get(ext, "不支持该文件格式"))

    if material_type == "UNSUPPORTED":
        raise ValueError(f"不支持的文件类型：{ext}")

    if material_type == "DOCX":
        if len(file_bytes) > MAX_NEW_TYPE_SIZE:
            raise ValueError(f"文件过大，当前类型最大支持 20MB，请压缩或拆分后上传。")
        text = extract_docx_text(file_bytes)
        return {"material_type": "DOCX", "text": text, "metadata": {"extension": ext, "filename": filename}}

    if material_type == "PPTX":
        if len(file_bytes) > MAX_NEW_TYPE_SIZE:
            raise ValueError(f"文件过大，当前类型最大支持 20MB，请压缩或拆分后上传。")
        text = extract_pptx_text(file_bytes)
        slide_count = text.count("【第") if text else 0
        return {
            "material_type": "PPTX",
            "text": text,
            "metadata": {"extension": ext, "filename": filename, "slides": slide_count},
        }

    if material_type == "TEXT":
        if len(file_bytes) > MAX_NEW_TYPE_SIZE:
            raise ValueError(f"文件过大，当前类型最大支持 20MB，请压缩或拆分后上传。")
        text = extract_text_file(file_bytes, filename)
        return {"material_type": "TEXT", "text": text, "metadata": {"extension": ext, "filename": filename}}

    if material_type == "CODE":
        if len(file_bytes) > MAX_NEW_TYPE_SIZE:
            raise ValueError(f"文件过大，当前类型最大支持 20MB，请压缩或拆分后上传。")
        text = extract_code_file(file_bytes, filename)
        language = CODE_LANGUAGE_MAP.get(ext, "Code")
        return {
            "material_type": "CODE",
            "text": text,
            "metadata": {"extension": ext, "filename": filename, "language": language},
        }

    raise ValueError(f"当前类型 {material_type} 暂不支持 extract_supported_file_text，请使用现有解析流程。")
